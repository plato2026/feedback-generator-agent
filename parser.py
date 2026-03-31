"""
parser.py — Enhanced Multi-Format Student Submission Parser
============================================================
Parses uploaded files to extract individual student submissions.
Supports: .docx, .txt, .pdf, .pptx, .xlsx, .png, .jpg, .jpeg, .zip

NEW FEATURES:
- PDF extraction
- PowerPoint extraction
- Excel extraction
- Image OCR (for charts/graphs)
- ZIP folder processing (Canvas exports)
- Multiple file uploads
"""

import re
import io
import zipfile
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass
from pathlib import Path


@dataclass
class StudentSubmission:
    """Represents a single parsed student submission."""
    name: str
    text: str
    word_count: int
    char_count: int
    source_file: str = "unknown"  # NEW: Track which file this came from
    file_type: str = "text"       # NEW: Track file type


# Delimiter patterns ordered from most specific to most general
DELIMITER_PATTERNS = [
    (r"^Student:\s*(.+)$",          "Student: [Name]"),
    (r"^Name:\s*(.+)$",             "Name: [Name]"),
    (r"^Submission by:?\s*(.+)$",   "Submission by: [Name]"),
    (r"^---\s*([A-Z][a-z]+ [A-Z].+?)\s*---$", "--- Name ---"),
    (r"^##\s*([A-Z][a-z]+ [A-Z].+)$",          "## Name"),
]


# ═══════════════════════════════════════════════════════════════════════
# TEXT EXTRACTION FUNCTIONS FOR DIFFERENT FILE TYPES
# ═══════════════════════════════════════════════════════════════════════

def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract plain text from a .docx file using mammoth."""
    try:
        import mammoth
        result = mammoth.extract_raw_text(io.BytesIO(file_bytes))
        return result.value
    except Exception:
        pass

    # Fallback: python-docx
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        text_parts = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text)
        
        return "\n\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Could not read .docx file: {e}")


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF file."""
    try:
        from pypdf import PdfReader
        pdf_reader = PdfReader(io.BytesIO(file_bytes))
        text_parts = []
        
        for page_num, page in enumerate(pdf_reader.pages, 1):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(f"--- Page {page_num} ---\n{page_text}")
        
        if not text_parts:
            return "⚠️ Could not extract text from PDF. The PDF might be image-based or encrypted."
        
        return "\n\n".join(text_parts)
    
    except Exception as e:
        return f"⚠️ Error reading PDF: {str(e)}"


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from a PowerPoint presentation."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extract table data
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_data:
                            slide_text.append(" | ".join(row_data))
            
            # Extract notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    slide_text.append(f"\n[Notes: {notes_text}]")
            
            if len(slide_text) > 1:  # More than just the header
                text_parts.append("\n".join(slide_text))
        
        if not text_parts:
            return "⚠️ Could not extract text from PowerPoint. Slides might be image-based."
        
        return "\n\n".join(text_parts)
    
    except Exception as e:
        return f"⚠️ Error reading PowerPoint: {str(e)}"


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    """Extract text from an Excel file."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
        text_parts = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_text = [f"--- Sheet: {sheet_name} ---"]
            
            for row in ws.iter_rows(values_only=True):
                # Filter out empty cells and convert to strings
                row_data = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if row_data:
                    sheet_text.append(" | ".join(row_data))
            
            if len(sheet_text) > 1:  # More than just the header
                text_parts.append("\n".join(sheet_text))
        
        if not text_parts:
            return "⚠️ Could not extract text from Excel. File might be empty."
        
        return "\n\n".join(text_parts)
    
    except Exception as e:
        return f"⚠️ Error reading Excel: {str(e)}"


def extract_text_from_image(file_bytes: bytes, filename: str) -> str:
    """
    Extract text from an image using OCR (pytesseract).
    Useful for charts, graphs, or scanned documents.
    """
    try:
        from PIL import Image
        import pytesseract
        
        image = Image.open(io.BytesIO(file_bytes))
        
        # Perform OCR
        text = pytesseract.image_to_string(image)
        
        if not text.strip():
            return f"⚠️ Could not extract text from image '{filename}'. Image might not contain readable text."
        
        return f"--- Image: {filename} ---\n{text}"
    
    except ImportError:
        return f"⚠️ OCR not available. Install pytesseract to extract text from images. Image: {filename}"
    except Exception as e:
        return f"⚠️ Error reading image '{filename}': {str(e)}"


def process_file_by_type(file_bytes: bytes, filename: str) -> Tuple[str, str]:
    """
    Process a file based on its extension and return extracted text and file type.
    
    Returns:
        Tuple of (extracted_text, file_type)
    """
    filename_lower = filename.lower()
    
    # Word documents
    if filename_lower.endswith(('.docx', '.doc')):
        return extract_text_from_docx(file_bytes), "docx"
    
    # PDF files
    elif filename_lower.endswith('.pdf'):
        return extract_text_from_pdf(file_bytes), "pdf"
    
    # PowerPoint files
    elif filename_lower.endswith(('.pptx', '.ppt')):
        return extract_text_from_pptx(file_bytes), "pptx"
    
    # Excel files
    elif filename_lower.endswith(('.xlsx', '.xls')):
        return extract_text_from_xlsx(file_bytes), "xlsx"
    
    # Image files
    elif filename_lower.endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
        return extract_text_from_image(file_bytes, filename), "image"
    
    # Plain text files
    elif filename_lower.endswith(('.txt', '.md', '.csv')):
        try:
            text = file_bytes.decode('utf-8')
            return text, "text"
        except UnicodeDecodeError:
            text = file_bytes.decode('latin-1')
            return text, "text"
    
    else:
        return f"⚠️ Unsupported file type: {filename}", "unknown"


# ═══════════════════════════════════════════════════════════════════════
# ZIP FILE HANDLING
# ═══════════════════════════════════════════════════════════════════════

def extract_files_from_zip(zip_bytes: bytes) -> List[Tuple[str, bytes, str]]:
    """
    Extract all files from a ZIP archive.
    
    Returns:
        List of tuples: (filename, file_bytes, relative_path)
    """
    extracted_files = []
    
    try:
        with zipfile.ZipFile(io.BytesIO(zip_bytes), 'r') as zip_ref:
            for file_info in zip_ref.filelist:
                # Skip directories and hidden files
                if file_info.is_dir() or file_info.filename.startswith('.'):
                    continue
                
                # Skip __MACOSX and other system folders
                if '__MACOSX' in file_info.filename or file_info.filename.startswith('_'):
                    continue
                
                # Extract file
                try:
                    file_bytes = zip_ref.read(file_info.filename)
                    filename = Path(file_info.filename).name  # Get just the filename
                    extracted_files.append((filename, file_bytes, file_info.filename))
                except Exception as e:
                    print(f"Warning: Could not extract {file_info.filename}: {e}")
                    continue
        
        return extracted_files
    
    except Exception as e:
        raise ValueError(f"Could not read ZIP file: {e}")


def process_multiple_files(files_data: List[Tuple[str, bytes]]) -> Tuple[str, Dict[str, int]]:
    """
    Process multiple files and combine their text.
    
    Args:
        files_data: List of (filename, file_bytes) tuples
    
    Returns:
        Tuple of (combined_text, file_stats)
    """
    all_text_parts = []
    file_stats = {
        'total': len(files_data),
        'docx': 0,
        'pdf': 0,
        'pptx': 0,
        'xlsx': 0,
        'image': 0,
        'text': 0,
        'unknown': 0,
        'errors': 0
    }
    
    for filename, file_bytes in files_data:
        try:
            text, file_type = process_file_by_type(file_bytes, filename)
            
            # Update stats
            if file_type in file_stats:
                file_stats[file_type] += 1
            
            # Add to combined text with clear delimiter
            all_text_parts.append(f"\n\n{'='*60}\n")
            all_text_parts.append(f"FILE: {filename}\n")
            all_text_parts.append(f"TYPE: {file_type}\n")
            all_text_parts.append(f"{'='*60}\n\n")
            all_text_parts.append(text)
            
            # Check for errors
            if text.startswith("⚠️"):
                file_stats['errors'] += 1
        
        except Exception as e:
            file_stats['errors'] += 1
            all_text_parts.append(f"\n\n⚠️ Error processing {filename}: {str(e)}\n\n")
    
    combined_text = "".join(all_text_parts)
    return combined_text, file_stats


# ═══════════════════════════════════════════════════════════════════════
# STUDENT PARSING (SAME AS BEFORE, BUT ENHANCED)
# ═══════════════════════════════════════════════════════════════════════

def parse_students(raw_text: str, source_file: str = "upload") -> Tuple[List[StudentSubmission], Optional[str], Optional[str]]:
    """
    Parse raw text to extract student submissions.
    
    Args:
        raw_text: The full text content
        source_file: Name of the source file (for tracking)
    
    Returns:
        Tuple of (list of StudentSubmission, error_message, pattern_used)
    """
    if not raw_text or len(raw_text.strip()) < 20:
        return [], "The uploaded content appears to be empty or too short to contain student submissions.", None

    # Try each delimiter pattern
    best_matches = []
    best_label = ""

    for pattern, label in DELIMITER_PATTERNS:
        matches = list(re.finditer(pattern, raw_text, re.MULTILINE | re.IGNORECASE))
        if len(matches) > len(best_matches):
            best_matches = matches
            best_label = label

    if not best_matches:
        return [], (
            "**Could not detect student delimiters.**\n\n"
            "Please ensure each student's submission is separated by a clear header line. "
            "The recommended format is:\n\n"
            "```\n"
            "Student: Jane Doe\n"
            "[Jane's complete submission text here...]\n\n"
            "Student: John Smith\n"
            "[John's complete submission text here...]\n"
            "```\n\n"
            "**For ZIP files with individual student files:**\n"
            "Each file should be named with the student's name, or contain a header with their name."
        ), None

    # Extract individual submissions
    students = []
    for i, match in enumerate(best_matches):
        name = match.group(1).strip().rstrip(":-_").strip()
        start_idx = match.end()
        end_idx = best_matches[i + 1].start() if i < len(best_matches) - 1 else len(raw_text)
        body = raw_text[start_idx:end_idx].strip()

        # Validate
        if 1 < len(name) < 80 and len(body) > 10:
            words = len(body.split())
            students.append(StudentSubmission(
                name=name,
                text=body,
                word_count=words,
                char_count=len(body),
                source_file=source_file
            ))

    if not students:
        return [], (
            "Delimiters were detected but no valid student submissions could be extracted. "
            "Please check that each student header is followed by submission text."
        ), best_label

    return students, None, best_label


# ═══════════════════════════════════════════════════════════════════════
# HELPER FUNCTION: PARSE FROM FILE NAME
# ═══════════════════════════════════════════════════════════════════════

def extract_student_name_from_filename(filename: str) -> Optional[str]:
    """
    Try to extract student name from filename.
    Common Canvas patterns:
    - "FirstName_LastName_123456_assignsubmission_file_..."
    - "LastName_FirstName_..."
    - "FirstName LastName..."
    """
    # Remove extension
    name_part = Path(filename).stem
    
    # Remove common Canvas patterns
    name_part = re.sub(r'_\d+_assignsubmission.*', '', name_part)
    name_part = re.sub(r'_submission.*', '', name_part, flags=re.IGNORECASE)
    name_part = re.sub(r'_\d{6,}', '', name_part)  # Remove student IDs
    
    # Replace underscores with spaces
    name_part = name_part.replace('_', ' ')
    
    # Basic validation: should have at least two parts (first and last name)
    parts = name_part.split()
    if len(parts) >= 2:
        return name_part.strip()
    
    return None


def create_submission_from_file(filename: str, text: str, file_type: str) -> StudentSubmission:
    """
    Create a StudentSubmission from a single file.
    Tries to extract student name from filename if not in text.
    """
    # Try to find student name in text
    for pattern, label in DELIMITER_PATTERNS:
        match = re.search(pattern, text, re.MULTILINE | re.IGNORECASE)
        if match:
            name = match.group(1).strip().rstrip(":-_").strip()
            # Remove the header from text
            text = re.sub(pattern, '', text, count=1, flags=re.MULTILINE | re.IGNORECASE).strip()
            break
    else:
        # Try filename
        name = extract_student_name_from_filename(filename)
        if not name:
            name = filename  # Fallback to filename
    
    words = len(text.split())
    
    return StudentSubmission(
        name=name,
        text=text,
        word_count=words,
        char_count=len(text),
        source_file=filename,
        file_type=file_type
    )
